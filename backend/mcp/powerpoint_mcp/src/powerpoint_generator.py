"""PowerPoint generator from JSON structure."""

import json
import os
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from loguru import logger
from src.schema import Presentation as PresentationSchema
from src.text_cleaner import clean_text


class PowerPointGenerator:
    """Generates PowerPoint presentations from JSON structure."""
    
    # Color scheme
    COLORS = {
        'primary': RGBColor(0, 51, 102),      # Dark blue
        'secondary': RGBColor(0, 102, 204),   # Light blue
        'accent': RGBColor(255, 140, 0),      # Orange
        'text': RGBColor(51, 51, 51),         # Dark gray
        'light': RGBColor(240, 240, 240),     # Light gray
        'white': RGBColor(255, 255, 255),     # White
        'success': RGBColor(46, 125, 50),     # Green
        'warning': RGBColor(255, 152, 0),     # Orange
    }
    
    def __init__(self):
        """Initialize PowerPoint generator."""
        self.prs = Presentation()
        self._setup_presentation()
        # Content width accounting for sidebar
        self.sidebar_width = Inches(0.6)
        self.content_width = Inches(10) - self.sidebar_width - Inches(0.1)  # Leave space for sidebar
        
        # Path to Foyer template image
        self.foyer_template_path = Path("template/foyer_template.jpg")
        if not self.foyer_template_path.exists():
            # Try alternative path relative to project root
            self.foyer_template_path = Path(__file__).parent.parent / "template" / "foyer_template.jpg"
        
        logger.info("PowerPoint generator initialized")
    
    def _setup_presentation(self):
        """Set up presentation settings."""
        # Set slide size to widescreen (16:9)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
    
    def generate_from_json(self, json_data: Dict[str, Any], output_path: Path) -> Path:
        """
        Generate PowerPoint from JSON data.
        
        Args:
            json_data: Parsed JSON presentation data
            output_path: Path to save the PowerPoint file
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Reset presentation for new generation
            self.prs = Presentation()
            self._setup_presentation()
            
            # Convert string to Path if necessary
            if isinstance(output_path, str):
                output_path = Path(output_path)
            
            # Validate JSON structure
            presentation = PresentationSchema.model_validate(json_data)
            
            # Generate slides
            for slide_data in presentation.slides:
                self._create_slide(slide_data.model_dump())
            
            # Save presentation
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(str(output_path))
            
            logger.success(f"PowerPoint generated: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to generate PowerPoint: {e}")
            raise
    
    def generate_from_json_file(self, json_path: Path, output_path: Optional[Path] = None) -> Path:
        """
        Generate PowerPoint from JSON file.
        
        Args:
            json_path: Path to JSON file
            output_path: Optional output path for PowerPoint
            
        Returns:
            Path to the generated PowerPoint file
        """
        if not json_path.exists():
            raise FileNotFoundError(f"JSON file not found: {json_path}")
        
        # Load JSON
        with open(json_path, 'r', encoding='utf-8') as f:
            json_data = json.load(f)
        
        # Generate default output path if not provided
        if output_path is None:
            output_path = json_path.parent / f"{json_path.stem}.pptx"
        
        return self.generate_from_json(json_data, output_path)
    
    def _create_slide(self, slide_data: Dict[str, Any]):
        """Create a slide based on layout type."""
        layout_type = slide_data.get('layout_type', 'bullet_points')
        
        # Map layout types to creation methods
        layout_handlers = {
            'title_slide': self._create_title_slide,
            'bullet_points': self._create_bullet_slide,
            'table': self._create_table_slide,
            'text_heavy': self._create_text_slide,
            'comparison': self._create_comparison_slide,
            'process_flow': self._create_process_slide,
            'mixed': self._create_mixed_slide,
        }
        
        handler = layout_handlers.get(layout_type, self._create_bullet_slide)
        handler(slide_data)
        
        # Add the Foyer sidebar to every slide
        self._add_foyer_sidebar(self.prs.slides[-1])
    
    def _create_title_slide(self, slide_data: Dict[str, Any]):
        """Create a title slide."""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = clean_text(slide_data.get('title', 'Presentation'))
        self._format_title(title)
        
        # Add subtitle if available
        if slide.placeholders[1]:
            subtitle = slide.placeholders[1]
            content = slide_data.get('content', {})
            subtitle.text = clean_text(content.get('subtitle', ''))
            self._format_subtitle(subtitle)
    
    def _create_bullet_slide(self, slide_data: Dict[str, Any]):
        """Create a bullet points slide."""
        slide_layout = self.prs.slide_layouts[1]  # Bullet slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = clean_text(slide_data.get('title', ''))
        self._format_title(title)
        
        # Add content
        content = slide_data.get('content', {})
        bullets = content.get('bullets', [])
        
        if bullets and slide.placeholders[1]:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()  # Clear default text
            
            # Add intro text if available
            if content.get('intro_text'):
                p = text_frame.add_paragraph()
                p.text = clean_text(content['intro_text'])
                p.level = 0
                self._format_paragraph(p)
                text_frame.add_paragraph()  # Add spacing
            
            # Add bullets
            for i, bullet in enumerate(bullets):
                if i > 0:  # Add paragraph for all except first
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                
                # Handle dict or string format
                if isinstance(bullet, dict):
                    p.text = clean_text(bullet.get('text', ''))
                    p.level = 0
                    self._format_paragraph(p)
                    
                    # Add sub-bullets
                    sub_bullets = bullet.get('sub_bullets', [])
                    if sub_bullets:  # Only iterate if sub_bullets is not None or empty
                        for sub in sub_bullets:
                            sp = text_frame.add_paragraph()
                            sp.text = clean_text(sub if isinstance(sub, str) else sub.get('text', ''))
                            sp.level = 1
                            self._format_paragraph(sp, is_sub=True)
                else:
                    p.text = clean_text(str(bullet))
                    p.level = 0
                    self._format_paragraph(p)
    
    def _create_table_slide(self, slide_data: Dict[str, Any]):
        """Create a table slide."""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), self.content_width - Inches(0.5), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = clean_text(slide_data.get('title', ''))
        self._format_title_text(title_frame.paragraphs[0])
        
        # Get table data
        content = slide_data.get('content', {})
        headers = content.get('headers', [])
        rows = content.get('rows', [])
        
        if not headers or not rows:
            return
        
        # Calculate table dimensions
        cols = len(headers)
        rows_count = len(rows) + 1  # +1 for header
        
        # Add table
        left = Inches(0.5)
        top = Inches(1.5)
        width = self.content_width - Inches(0.5)
        height = Inches(0.5 * min(rows_count, 8))  # Limit height
        
        table = slide.shapes.add_table(rows_count, cols, left, top, width, height).table
        
        # Set header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = clean_text(str(header))
            self._format_table_cell(cell, is_header=True)
        
        # Add data rows
        for row_idx, row_data in enumerate(rows, 1):
            for col_idx, cell_data in enumerate(row_data[:cols]):  # Ensure we don't exceed columns
                cell = table.cell(row_idx, col_idx)
                cell.text = clean_text(str(cell_data))
                self._format_table_cell(cell)
        
        # Add caption if available
        if content.get('caption'):
            caption_box = slide.shapes.add_textbox(
                Inches(0.5), top + height + Inches(0.2), self.content_width - Inches(0.5), Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = clean_text(content['caption'])
            self._format_caption(caption_frame.paragraphs[0])
    
    def _create_text_slide(self, slide_data: Dict[str, Any]):
        """Create a text-heavy slide."""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), self.content_width - Inches(0.5), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = clean_text(slide_data.get('title', ''))
        self._format_title_text(title_frame.paragraphs[0])
        
        # Add content
        content = slide_data.get('content', {})
        
        # Create main text box
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), self.content_width - Inches(0.5), Inches(3.8)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add paragraphs
        paragraphs = content.get('paragraphs', [])
        for i, para in enumerate(paragraphs):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = clean_text(para)
            self._format_paragraph(p, size=14)
            
            if i < len(paragraphs) - 1:
                text_frame.add_paragraph()  # Add spacing
        
        # Add key points if available
        if content.get('key_points'):
            text_frame.add_paragraph()
            p = text_frame.add_paragraph()
            p.text = "Points clés:"
            self._format_paragraph(p, bold=True)
            
            for point in content['key_points']:
                p = text_frame.add_paragraph()
                p.text = clean_text(f"• {point}")
                self._format_paragraph(p, size=12)
        
        # Add emphasis if available
        if content.get('emphasis'):
            text_frame.add_paragraph()
            p = text_frame.add_paragraph()
            p.text = clean_text(content['emphasis'])
            self._format_paragraph(p, bold=True, color=self.COLORS['accent'])
    
    def _create_comparison_slide(self, slide_data: Dict[str, Any]):
        """Create a comparison slide."""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), self.content_width - Inches(0.5), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = clean_text(slide_data.get('title', ''))
        self._format_title_text(title_frame.paragraphs[0])
        
        # Get comparison data
        content = slide_data.get('content', {})
        left_data = content.get('left', {})
        right_data = content.get('right', {})
        
        # Add comparison title if available
        if content.get('comparison_title'):
            comp_title_box = slide.shapes.add_textbox(
                Inches(2), Inches(1.2), Inches(5.5), Inches(0.4)
            )
            comp_title_frame = comp_title_box.text_frame
            comp_title_frame.text = clean_text(content['comparison_title'])
            p = comp_title_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._format_paragraph(p, bold=True, size=16)
        
        # Create left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(4.2), Inches(3.2)
        )
        left_frame = left_box.text_frame
        
        # Add left title
        p = left_frame.paragraphs[0]
        p.text = clean_text(left_data.get('title', 'Option A'))
        p.alignment = PP_ALIGN.CENTER
        self._format_paragraph(p, bold=True, color=self.COLORS['primary'])
        
        # Add left points
        for point in left_data.get('points', []):
            p = left_frame.add_paragraph()
            p.text = clean_text(f"• {point}")
            self._format_paragraph(p, size=12)
        
        # Create right column
        right_box = slide.shapes.add_textbox(
            Inches(5.3), Inches(1.8), Inches(4.2), Inches(3.2)
        )
        right_frame = right_box.text_frame
        
        # Add right title
        p = right_frame.paragraphs[0]
        p.text = clean_text(right_data.get('title', 'Option B'))
        p.alignment = PP_ALIGN.CENTER
        self._format_paragraph(p, bold=True, color=self.COLORS['secondary'])
        
        # Add right points
        for point in right_data.get('points', []):
            p = right_frame.add_paragraph()
            p.text = clean_text(f"• {point}")
            self._format_paragraph(p, size=12)
        
        # Add visual separator
        line = slide.shapes.add_connector(
            1, Inches(4.95), Inches(1.8), Inches(4.95), Inches(5)
        )
        line.line.color.rgb = self.COLORS['light']
        line.line.width = Pt(2)
    
    def _create_process_slide(self, slide_data: Dict[str, Any]):
        """Create a process flow slide."""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), self.content_width - Inches(0.5), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = clean_text(slide_data.get('title', ''))
        self._format_title_text(title_frame.paragraphs[0])
        
        # Get process data
        content = slide_data.get('content', {})
        steps = content.get('steps', [])
        flow_type = content.get('flow_type', 'linear')
        
        if not steps:
            return
        
        # Sort steps by order
        steps = sorted(steps, key=lambda x: x.get('order', 0))
        
        # Calculate positions
        num_steps = len(steps)
        if num_steps <= 4:
            # Horizontal layout
            step_width = Inches(2)
            step_height = Inches(1.2)
            spacing = Inches(0.3)
            total_width = num_steps * step_width + (num_steps - 1) * spacing
            start_x = (Inches(10) - total_width) / 2
            y_pos = Inches(2.5)
            
            for i, step in enumerate(steps):
                x_pos = start_x + i * (step_width + spacing)
                
                # Add step box
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x_pos, y_pos, step_width, step_height
                )
                
                # Style the shape
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.COLORS['primary'] if i % 2 == 0 else self.COLORS['secondary']
                shape.line.color.rgb = self.COLORS['white']
                
                # Add text
                text_frame = shape.text_frame
                text_frame.clear()
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                # Add step title
                p = text_frame.add_paragraph()
                p.text = clean_text(f"{i+1}. {step.get('title', '')}")
                p.alignment = PP_ALIGN.CENTER
                self._format_paragraph(p, bold=True, color=self.COLORS['white'], size=14)
                
                # Add description if available
                if step.get('description'):
                    p = text_frame.add_paragraph()
                    p.text = clean_text(step['description'])
                    p.alignment = PP_ALIGN.CENTER
                    self._format_paragraph(p, color=self.COLORS['white'], size=10)
                
                # Add arrow (except for last step)
                if i < num_steps - 1:
                    arrow_x = x_pos + step_width + Inches(0.05)
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, y_pos + step_height/2 - Inches(0.2),
                        Inches(0.2), Inches(0.4)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self.COLORS['accent']
        else:
            # Vertical or grid layout for more steps
            cols = 2
            rows = (num_steps + 1) // 2
            step_width = Inches(4)
            step_height = Inches(0.8)
            x_spacing = Inches(0.5)
            y_spacing = Inches(0.2)
            
            for i, step in enumerate(steps):
                col = i % cols
                row = i // cols
                
                x_pos = Inches(0.5) + col * (step_width + x_spacing)
                y_pos = Inches(1.5) + row * (step_height + y_spacing)
                
                # Add step box
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x_pos, y_pos, step_width, step_height
                )
                
                # Style the shape
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.COLORS['primary'] if i % 2 == 0 else self.COLORS['secondary']
                
                # Add text
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.add_paragraph()
                title_text = clean_text(f"{i+1}. {step.get('title', '')}")
                if step.get('description'):
                    title_text += f": {clean_text(step['description'])}"
                p.text = title_text
                p.alignment = PP_ALIGN.LEFT
                self._format_paragraph(p, color=self.COLORS['white'], size=11)
    
    def _create_mixed_slide(self, slide_data: Dict[str, Any]):
        """Create a mixed content slide."""
        # For now, treat as bullet slide
        # Can be enhanced to handle mixed content more sophisticatedly
        self._create_bullet_slide(slide_data)
    
    # Formatting helper methods
    def _format_title(self, title_placeholder):
        """Format title placeholder."""
        if title_placeholder and title_placeholder.text_frame:
            for paragraph in title_placeholder.text_frame.paragraphs:
                self._format_title_text(paragraph)
    
    def _format_subtitle(self, subtitle_placeholder):
        """Format subtitle placeholder."""
        if subtitle_placeholder and subtitle_placeholder.text_frame:
            for paragraph in subtitle_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = self.COLORS['text']
    
    def _format_title_text(self, paragraph):
        """Format title text."""
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.COLORS['primary']
        paragraph.alignment = PP_ALIGN.LEFT
    
    def _format_paragraph(self, paragraph, size=16, bold=False, is_sub=False, color=None):
        """Format paragraph text."""
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color or self.COLORS['text']
        
        if is_sub:
            paragraph.font.size = Pt(14)
    
    def _format_table_cell(self, cell, is_header=False):
        """Format table cell."""
        if cell.text_frame and cell.text_frame.paragraphs:
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12 if not is_header else 14)
            p.font.bold = is_header
            
            if is_header:
                p.font.color.rgb = self.COLORS['white']
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.COLORS['primary']
            else:
                p.font.color.rgb = self.COLORS['text']
    
    def _format_caption(self, paragraph):
        """Format caption text."""
        paragraph.font.size = Pt(10)
        paragraph.font.italic = True
        paragraph.font.color.rgb = self.COLORS['text']
        paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_foyer_sidebar(self, slide):
        """Add Foyer sidebar image to the slide."""
        try:
            # Check if template image exists
            if not self.foyer_template_path.exists():
                logger.warning(f"Foyer template not found at {self.foyer_template_path}")
                # Fall back to creating a simple blue bar if image not found
                self._add_fallback_sidebar(slide)
                return
            
            # Get image dimensions to maintain aspect ratio
            from PIL import Image
            with Image.open(self.foyer_template_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
            
            # Calculate proper dimensions maintaining aspect ratio
            # Use the full slide height and calculate width based on aspect ratio
            desired_height = self.prs.slide_height
            desired_width = desired_height * aspect_ratio
            
            # Position image on the right edge
            img_left = self.prs.slide_width - desired_width
            img_top = Inches(0)
            
            # Add the image with proper aspect ratio
            pic = slide.shapes.add_picture(
                str(self.foyer_template_path),
                img_left,
                img_top,
                width=desired_width,
                height=desired_height
            )
            
            # Lock the image to prevent editing/moving
            # Make it behave like a background element
            pic.is_decorative = True  # Mark as decorative/background element
            
            # Try to lock the shape (this might not work in all PowerPoint versions)
            try:
                # Access the shape's lock properties
                shape = pic._element
                shape.set('locked', '1')
            except:
                pass  # If locking doesn't work, continue anyway
            
            # Update content width based on actual image width
            self.sidebar_width = desired_width
            self.content_width = self.prs.slide_width - desired_width - Inches(0.1)
            
            logger.debug(f"Added Foyer sidebar image from {self.foyer_template_path} with aspect ratio {aspect_ratio:.2f}")
            
        except Exception as e:
            logger.error(f"Failed to add Foyer template image: {e}")
            # Fall back to simple sidebar
            self._add_fallback_sidebar(slide)
    
    def _add_fallback_sidebar(self, slide):
        """Add a simple blue sidebar as fallback when image is not available."""
        # Bar dimensions
        bar_width = self.sidebar_width
        bar_height = self.prs.slide_height
        bar_left = self.prs.slide_width - bar_width
        bar_top = Inches(0)
        
        # Create simple blue rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left, bar_top, bar_width, bar_height
        )
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 115, 190)  # Blue
        shape.line.fill.background()  # No border
        
        # Add Foyer text
        text_box = slide.shapes.add_textbox(
            bar_left,
            bar_height - Inches(1),
            bar_width,
            Inches(0.5)
        )
        
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = "Foyer"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)